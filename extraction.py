import os
import zipfile
import uuid
import io
import re

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
from crud import save_file  # Импортируем save_file для добавления файлов в базу

def docx_paragraph_label(p):
    style = p.style.name.lower() if p.style else ""
    if p._p.pPr is not None and p._p.pPr.numPr is not None:
        return "LIST ITEM"
    if "heading" in style:
        return "HEADING"
    for run in p.runs:
        if run.bold:
            return "HEADING"
        if run.italic:
            return "SUBHEADING"
    return "NORMAL"


def pptx_paragraph_label(p):
    if p.level and p.level > 0:
        return "LIST ITEM"
    bold = any(r.font.bold for r in p.runs)
    italic = any(r.font.italic for r in p.runs)
    sizes = [r.font.size.pt for r in p.runs if r.font.size]
    if bold or (sizes and max(sizes) > 20):
        return "HEADING"
    if italic:
        return "SUBHEADING"
    return "NORMAL"


def markdown_from_label(label, txt: str):
    if label == "HEADING":
        return f"## {txt}"
    if label == "SUBHEADING":
        return f"### {txt}"
    if label == "LIST ITEM":
        return f"- {txt}"
    return txt


def extract_from_docx(path, raw_lines, images_dir, project_id, db):
    doc = Document(path)
    raw_lines.append(f"---\n# {os.path.basename(path)}\n---")
    extracted_images = []
    for p in doc.paragraphs:
        txt = p.text.strip()
        if txt:
            lbl = docx_paragraph_label(p)
            raw_lines.append(markdown_from_label(lbl, txt))
            raw_lines.append("")
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            raw_lines.append("| " + " | ".join(cells) + " |")
    with zipfile.ZipFile(path, "r") as z:
        for info in z.infolist():
            if info.filename.startswith("word/media/"):
                data = z.read(info.filename)
                fn = f"{uuid.uuid4()}_{os.path.basename(info.filename)}"
                open(os.path.join(images_dir, fn), "wb").write(data)
                extracted_images.append(fn)
                # Регистрируем извлеченный файл в базе
                save_file(db, project_id, filename=fn, file_type="image/jpeg")
    return extracted_images


def extract_from_pptx(path, raw_lines, images_dir, project_id, db):
    prs = Presentation(path)
    raw_lines.append(f"---\n# {os.path.basename(path)}\n---")
    extracted_images = []
    for i, slide in enumerate(prs.slides, start=1):
        raw_lines.append(f"## Slide {i}")
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    lbl = pptx_paragraph_label(p)
                    raw_lines.append(markdown_from_label(lbl, txt))
        for shp in slide.shapes:
            if shp.shape_type == 13 and shp.image:
                blob = shp.image.blob
                fn = f"slide{i}_{uuid.uuid4()}.png"
                Image.open(io.BytesIO(blob)).save(os.path.join(images_dir, fn))
                extracted_images.append(fn)
                # Регистрируем извлеченный файл в базе
                save_file(db, project_id, filename=fn, file_type="image/png")
    return extracted_images


def extract_from_pdf(path, raw_lines, images_dir, project_id, db):
    raw_lines.append(f"---\n# {os.path.basename(path)}\n---")
    extracted_images = []
    doc = fitz.open(path)
    for i, page in enumerate(doc, start=1):
        txt = page.get_text().strip()
        if txt:
            raw_lines.append(f"## Page {i}")
            for line in txt.splitlines():
                raw_lines.append(line)
        for idx, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.n < 5:
                mode = "RGB" if pix.n == 3 else "L"
                image = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
            else:
                pix = fitz.Pixmap(fitz.csRGB, pix)
                image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            fn = f"page{i}_img{idx}_{uuid.uuid4()}.png"
            image.save(os.path.join(images_dir, fn))
            extracted_images.append(fn)
            # Регистрируем извлеченный файл в базе
            save_file(db, project_id, filename=fn, file_type="image/png")
            pix = None
    doc.close()
    return extracted_images


def extract_text_and_images(file_paths, project_dir, project_id, db):
    raw = []
    img_dir = os.path.join(project_dir, "images")
    os.makedirs(img_dir, exist_ok=True)

    for p in file_paths:
        ext = os.path.splitext(p)[1].lower().lstrip('.')
        func = globals().get(f"extract_from_{ext}")
        if callable(func):
            func(p, raw, img_dir, project_id, db)

    return "\n".join(raw)